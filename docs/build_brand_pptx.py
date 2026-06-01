#!/usr/bin/env python3
"""Build the Run 2 The Rescue Brand Playbook .pptx.

Output: docs/RUN_2_THE_RESCUE_BRAND_PLAYBOOK.pptx

Run from the site/ directory:
    pip install python-pptx
    python3 docs/build_brand_pptx.py

Rebuild any time the brand evolves; tokens here mirror styles.css.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---- Brand colors (from styles.css :root) ----
PLUM_900    = RGBColor(0x17, 0x10, 0x25)
PLUM_800    = RGBColor(0x1f, 0x15, 0x30)
PLUM_700    = RGBColor(0x2a, 0x1c, 0x3d)
PLUM_600    = RGBColor(0x35, 0x24, 0x49)
PLUM_500    = RGBColor(0x43, 0x32, 0x5a)
PURPLE_400  = RGBColor(0xb4, 0x8b, 0xdf)
PURPLE_500  = RGBColor(0x98, 0x71, 0xd2)
PURPLE_600  = RGBColor(0x7d, 0x56, 0xb9)
PURPLE_700  = RGBColor(0x5e, 0x3f, 0x95)
PURPLE_SOFT = RGBColor(0xe6, 0xd8, 0xf0)
LAV_50      = RGBColor(0xf5, 0xf3, 0xf7)
LAV_100     = RGBColor(0xeb, 0xe7, 0xef)
LAV_200     = RGBColor(0xdd, 0xd6, 0xe5)
LAV_300     = RGBColor(0xc7, 0xbc, 0xd5)
INK         = RGBColor(0x1a, 0x12, 0x26)
INK_2       = RGBColor(0x43, 0x39, 0x4e)
INK_3       = RGBColor(0x7b, 0x75, 0x85)
WHITE       = RGBColor(0xff, 0xff, 0xff)
URGENCY     = RGBColor(0xd9, 0x68, 0x47)
NEAR_WHITE  = RGBColor(0xcc, 0xc6, 0xd5)

# ---- Fonts (system fonts that look professional on Mac + Windows) ----
F_DISPLAY = "Helvetica Neue"
F_BODY    = "Helvetica Neue"
F_MONO    = "Menlo"

# ---- Slide geometry: 16:9 widescreen ----
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)


# ---- Helpers ----

def set_slide_bg(slide, color):
    """Fill the whole slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line=False):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height, fill_color, radius=0.08):
    """Add a rounded rectangle. Adjust radius via shape adjustment."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text,
             font=F_BODY, size=14, color=INK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """Add a text box with one paragraph."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box, tf, p


def add_paragraphs(slide, left, top, width, height,
                   paragraphs, font=F_BODY, size=12, color=INK_2,
                   line_spacing=1.4, space_after_pt=6,
                   align=PP_ALIGN.LEFT):
    """Add a text box with multiple paragraphs.
    paragraphs: list of (text, {opts}) tuples. opts may override font/size/color/bold/italic.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    for i, item in enumerate(paragraphs):
        text, opts = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        p.space_after = Pt(opts.get("space_after", space_after_pt))
        run = p.add_run()
        run.text = text
        run.font.name = opts.get("font", font)
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", color)
    return box


def page_chrome(slide, page_num, total, dark=False):
    """Tiny header label + page indicator + tagline on every content slide."""
    fg = LAV_200 if dark else INK_3
    # Top accent bar
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05),
                   PURPLE_500)
    # Header label
    add_text(slide, MARGIN, Inches(0.25), Inches(8), Inches(0.3),
             "RUN 2 THE RESCUE   ·   BRAND PLAYBOOK",
             font=F_MONO, size=8, color=fg)
    # Page indicator
    add_text(slide, SLIDE_W - MARGIN - Inches(2), Inches(0.25),
             Inches(2), Inches(0.3),
             f"{page_num:02d} / {total:02d}",
             font=F_MONO, size=8, color=fg, align=PP_ALIGN.RIGHT)


def eyebrow(slide, left, top, text, color=PURPLE_500):
    add_text(slide, left, top, Inches(8), Inches(0.3),
             f"✦   {text}",
             font=F_MONO, size=10, color=color, bold=True)


def slide_title(slide, text, color=INK, top=Inches(1.15), size=40):
    # Taller box so 2-line wrapping doesn't overlap the next element.
    add_text(slide, MARGIN, top, Inches(12), Inches(1.8),
             text, font=F_DISPLAY, size=size, color=color, bold=True,
             line_spacing=1.08)


# ---- Custom slide builders ----

def make_cover(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PLUM_900)
    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, PURPLE_500)
    # Eyebrow
    add_text(slide, MARGIN, Inches(0.9), Inches(8), Inches(0.4),
             "✦   RUN 2 THE RESCUE",
             font=F_MONO, size=11, color=PURPLE_400, bold=True)
    # Title
    add_text(slide, MARGIN, Inches(2.0), Inches(11), Inches(1.5),
             "Brand", font=F_DISPLAY, size=120, color=WHITE, bold=True,
             line_spacing=0.95)
    add_text(slide, MARGIN, Inches(3.3), Inches(11), Inches(1.5),
             "Playbook", font=F_DISPLAY, size=120, color=WHITE, bold=True,
             line_spacing=0.95)
    # Subtitle
    add_text(slide, MARGIN, Inches(5.0), Inches(11), Inches(0.45),
             "Voice, typography, color, and motion for",
             font=F_BODY, size=18, color=NEAR_WHITE)
    add_text(slide, MARGIN, Inches(5.4), Inches(11), Inches(0.45),
             "a nonprofit rescuing dogs from the East Asia meat trade.",
             font=F_BODY, size=18, color=NEAR_WHITE)
    # Divider
    add_rect(slide, MARGIN, Inches(6.5), Inches(4.5), Inches(0.02),
             PURPLE_700)
    # Tagline
    add_text(slide, MARGIN, Inches(6.65), Inches(8), Inches(0.4),
             "RUN.   RESCUE.   REPEAT.",
             font=F_DISPLAY, size=20, color=PURPLE_400, bold=True)
    # Meta
    add_text(slide, MARGIN, Inches(7.05), Inches(8), Inches(0.3),
             "501(c)(3) Nonprofit   ·   EIN 99-4240461   ·   Edition 2026",
             font=F_MONO, size=9, color=RGBColor(0xa8, 0x9d, 0xb4))


def make_foreword(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "FOREWORD")
    slide_title(slide, "How to use this playbook.")
    add_text(slide, MARGIN, Inches(2.8), Inches(11.5), Inches(0.9),
             "Every color, type choice, and voice rule here is already in use at",
             font=F_BODY, size=18, color=INK_2)
    add_text(slide, MARGIN, Inches(3.15), Inches(11.5), Inches(0.5),
             "run2-rescuedemo.vercel.app",
             font=F_MONO, size=18, color=PURPLE_700, bold=True)
    add_text(slide, MARGIN, Inches(3.55), Inches(11.5), Inches(0.5),
             "Open the live site in another tab while you read.",
             font=F_BODY, size=18, color=INK_2)
    # Two rules
    add_rect(slide, MARGIN, Inches(4.7), Inches(12), Inches(0.02), LAV_200)
    add_text(slide, MARGIN, Inches(4.9), Inches(12), Inches(0.4),
             "Two rules never break:",
             font=F_BODY, size=14, color=INK_3, bold=True)
    add_text(slide, MARGIN, Inches(5.35), Inches(12), Inches(0.5),
             "1.   The protected lexicon  (§ 03)",
             font=F_BODY, size=16, color=INK, bold=True)
    add_text(slide, MARGIN, Inches(5.75), Inches(12), Inches(0.5),
             "2.   No graphic cruelty outside the Reality page  (§ 07)",
             font=F_BODY, size=16, color=INK, bold=True)
    add_text(slide, MARGIN, Inches(6.25), Inches(12), Inches(0.4),
             "Everything else is guidance. Ship the better thing.",
             font=F_BODY, size=14, color=INK_3, italic=True)


def make_who_we_are(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "01   WHO WE ARE")
    slide_title(slide, "Dogs from the meat trade.")
    add_text(slide, MARGIN, Inches(2.85), Inches(12), Inches(0.6),
             "We pull them. We get them well. We send them home.",
             font=F_DISPLAY, size=28, color=INK_2, italic=True)
    # Facts grid
    facts = [
        ("FOUNDED", "2012", "Brandy Cherven (CEO) and Bonnie Klapper (COO)"),
        ("REGIONS", "China  ·  South Korea", "On-the-ground rescuer partners"),
        ("FLIGHT HUBS", "JFK  ·  LAX", "Where survivors land in the U.S."),
        ("STATUS", "501(c)(3)", "EIN 99-4240461  ·  Tax deductible"),
    ]
    col_w = Inches(2.85)
    gap = Inches(0.25)
    start_x = MARGIN
    y = Inches(4.2)
    for i, (label, big, sub) in enumerate(facts):
        x = start_x + (col_w + gap) * i
        add_text(slide, x, y, col_w, Inches(0.3), label,
                 font=F_MONO, size=9, color=PURPLE_600, bold=True)
        add_text(slide, x, y + Inches(0.35), col_w, Inches(0.8), big,
                 font=F_DISPLAY, size=24, color=INK, bold=True,
                 line_spacing=1.05)
        add_text(slide, x, y + Inches(1.25), col_w, Inches(0.9), sub,
                 font=F_BODY, size=11, color=INK_3, line_spacing=1.35)


def make_voice_attributes(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "02   VOICE AND TONE")
    slide_title(slide, "Four voices, blended.")
    rows = [
        ("Hopeful",          "We lead with what's possible, not what's painful."),
        ("Dignified",        "We treat both the dogs and the reader as worthy."),
        ("Urgent",           "Real timelines, real stakes. Never alarmist."),
        ("Never exploitative", "We never use trauma porn or guilt to drive action."),
    ]
    y = Inches(2.8)
    row_h = Inches(0.95)
    for label, desc in rows:
        add_rect(slide, MARGIN, y, Inches(0.04), row_h - Inches(0.15),
                 PURPLE_500)
        add_text(slide, MARGIN + Inches(0.25), y, Inches(3.2), Inches(0.5),
                 label, font=F_DISPLAY, size=22, color=INK, bold=True)
        add_text(slide, MARGIN + Inches(0.25), y + Inches(0.42),
                 Inches(11), Inches(0.5),
                 desc, font=F_BODY, size=14, color=INK_2)
        y += row_h


def make_reader_posture(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "02   READER POSTURE")
    slide_title(slide, "Every CTA is a door, not an ask.")
    add_text(slide, MARGIN, Inches(2.3), Inches(11), Inches(0.5),
             "The visitor is a co-rescuer, not a savior. Phrasing reflects it.",
             font=F_BODY, size=16, color=INK_3, italic=True)
    # Two columns: Use / Avoid
    col_w = Inches(5.8)
    use_x = MARGIN
    avoid_x = MARGIN + col_w + Inches(0.4)
    y_top = Inches(3.4)
    add_text(slide, use_x, y_top, col_w, Inches(0.3), "USE",
             font=F_MONO, size=10, color=PURPLE_700, bold=True)
    add_text(slide, avoid_x, y_top, col_w, Inches(0.3), "AVOID",
             font=F_MONO, size=10, color=URGENCY, bold=True)
    # Underlines
    add_rect(slide, use_x, y_top + Inches(0.32), col_w, Inches(0.02),
             PURPLE_500)
    add_rect(slide, avoid_x, y_top + Inches(0.32), col_w, Inches(0.02),
             URGENCY)
    pairs = [
        ("Meet the survivors",  "Please help us"),
        ("Bring one home",      "Donate to save lives"),
        ("Fund a flight",       "We need your support"),
        ("Open your home",      "Will you give?"),
    ]
    y = y_top + Inches(0.55)
    for use, avoid in pairs:
        add_text(slide, use_x, y, col_w, Inches(0.4), use,
                 font=F_BODY, size=18, color=INK, bold=True)
        add_text(slide, avoid_x, y, col_w, Inches(0.4), avoid,
                 font=F_BODY, size=18, color=INK_3, italic=True)
        y += Inches(0.55)


def make_transformation(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PLUM_900)
    page_chrome(slide, page, total, dark=True)
    eyebrow(slide, MARGIN, Inches(0.8), "02   TRANSFORMATION LANGUAGE",
            color=PURPLE_400)
    slide_title(slide, "Every survivor's arc bends through one of these.",
                color=WHITE)
    add_text(slide, MARGIN, Inches(3.3), Inches(12), Inches(1.1),
             "trauma to trust",
             font=F_DISPLAY, size=56, color=PURPLE_400, bold=True,
             line_spacing=1.05)
    add_text(slide, MARGIN, Inches(4.4), Inches(12), Inches(1.1),
             "fear to faith",
             font=F_DISPLAY, size=56, color=PURPLE_400, bold=True,
             line_spacing=1.05)
    add_text(slide, MARGIN, Inches(5.5), Inches(12), Inches(1.1),
             "forgotten to forever",
             font=F_DISPLAY, size=56, color=PURPLE_400, bold=True,
             line_spacing=1.05)
    add_text(slide, MARGIN, Inches(6.8), Inches(12), Inches(0.4),
             "Use sparingly. Never all three in one paragraph.",
             font=F_BODY, size=14, color=NEAR_WHITE, italic=True)


def make_protected_lexicon(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "03   THE PROTECTED LEXICON")
    slide_title(slide, "Brand assets. Never rewrite.")
    items = [
        ("RUN. RESCUE. REPEAT.",      "The tagline. Punctuation matters."),
        ('"Run to the rescue with love, and peace will follow."',
         "Anchor quote. River Phoenix."),
        ("hopeful survivors",         "How we refer to dogs in our care."),
        ("bravehearted survivors",    "Emphatic variant."),
        ("voiceless victims",         "Dogs still inside the trade."),
        ("second chance",             "What adoption represents."),
        ("beacon of hope",            "What R2TR is to activists and dogs."),
        ("Sponsor Angel",             "Monthly sponsorship program name."),
    ]
    y = Inches(2.7)
    for phrase, note in items:
        add_text(slide, MARGIN, y, Inches(7), Inches(0.36), phrase,
                 font=F_DISPLAY, size=15, color=INK, bold=True)
        add_text(slide, MARGIN + Inches(7.2), y, Inches(5), Inches(0.36),
                 note, font=F_BODY, size=12, color=INK_3, italic=True)
        y += Inches(0.45)


def color_swatch(slide, left, top, width, height,
                 token, oklch, hex_str, fill_color, light_text=True):
    add_round_rect(slide, left, top, width, height, fill_color, radius=0.06)
    tc = WHITE if light_text else INK
    add_text(slide, left + Inches(0.15), top + Inches(0.12),
             width - Inches(0.3), Inches(0.4),
             token, font=F_DISPLAY, size=14, color=tc, bold=True)
    add_text(slide, left + Inches(0.15), top + Inches(0.55),
             width - Inches(0.3), Inches(0.3),
             oklch, font=F_MONO, size=9, color=tc)
    add_text(slide, left + Inches(0.15), top + Inches(0.85),
             width - Inches(0.3), Inches(0.3),
             hex_str, font=F_MONO, size=9, color=tc)


def make_color_plum(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "05   COLOR  ·  PLUM")
    slide_title(slide, "Dark surfaces. The quiet base.")
    swatches = [
        ("PLUM 900", "oklch(0.18 0.035 310)", "#171025", PLUM_900),
        ("PLUM 800", "oklch(0.22 0.04 310)",  "#1f1530", PLUM_800),
        ("PLUM 700", "oklch(0.28 0.045 310)", "#2a1c3d", PLUM_700),
        ("PLUM 600", "oklch(0.34 0.05 310)",  "#352449", PLUM_600),
        ("PLUM 500", "oklch(0.42 0.055 310)", "#43325a", PLUM_500),
    ]
    w = Inches(2.36)
    gap = Inches(0.15)
    h = Inches(1.5)
    y = Inches(3.4)
    for i, (token, ok, hx, color) in enumerate(swatches):
        x = MARGIN + (w + gap) * i
        color_swatch(slide, x, y, w, h, token, ok, hx, color)
    add_text(slide, MARGIN, Inches(5.4), Inches(12), Inches(0.4),
             "Used for: page backgrounds on dark sections  ·  card backgrounds  ·  borders.",
             font=F_BODY, size=13, color=INK_3, italic=True)


def make_color_purple(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "05   COLOR  ·  PURPLE")
    slide_title(slide, "Accent and interaction.")
    swatches = [
        ("PURPLE SOFT", "oklch(0.92 0.05 305)", "#e6d8f0", PURPLE_SOFT, False),
        ("PURPLE 400", "oklch(0.72 0.14 305)", "#b48bdf", PURPLE_400, True),
        ("PURPLE 500", "oklch(0.63 0.16 305)", "#9871d2", PURPLE_500, True),
        ("PURPLE 600", "oklch(0.54 0.17 305)", "#7d56b9", PURPLE_600, True),
        ("PURPLE 700", "oklch(0.44 0.15 305)", "#5e3f95", PURPLE_700, True),
    ]
    w = Inches(2.36)
    gap = Inches(0.15)
    h = Inches(1.5)
    y = Inches(3.4)
    for i, (token, ok, hx, color, lt) in enumerate(swatches):
        x = MARGIN + (w + gap) * i
        color_swatch(slide, x, y, w, h, token, ok, hx, color, light_text=lt)
    add_text(slide, MARGIN, Inches(5.4), Inches(12), Inches(0.4),
             "PURPLE 500 is the only color that carries primary actions. One per screen.",
             font=F_BODY, size=13, color=INK_3, italic=True)


def make_color_lav_ink(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "05   COLOR  ·  LAVENDER + INK + URGENCY")
    slide_title(slide, "Light. Type. One accent.")
    # Lavender row
    lav = [
        ("LAV 50",  "oklch(0.97 0.012 300)", "#f5f3f7", LAV_50),
        ("LAV 100", "oklch(0.94 0.02 300)",  "#ebe7ef", LAV_100),
        ("LAV 200", "oklch(0.90 0.028 300)", "#ddd6e5", LAV_200),
        ("LAV 300", "oklch(0.84 0.035 300)", "#c7bcd5", LAV_300),
    ]
    w = Inches(2.6); gap = Inches(0.2); h = Inches(1.35)
    y1 = Inches(3.2)
    for i, (token, ok, hx, color) in enumerate(lav):
        x = MARGIN + (w + gap) * i
        color_swatch(slide, x, y1, w, h, token, ok, hx, color, light_text=False)
    # Ink row
    add_text(slide, MARGIN, Inches(4.85), Inches(8), Inches(0.3), "TEXT ON LIGHT",
             font=F_MONO, size=9, color=PURPLE_600, bold=True)
    inks = [
        ("INK",   "oklch(0.16 0.03 310)", "#1a1226", INK),
        ("INK 2", "oklch(0.36 0.03 310)", "#43394e", INK_2),
        ("INK 3", "oklch(0.55 0.02 310)", "#7b7585", INK_3),
    ]
    y2 = Inches(5.25)
    iw = Inches(2.6); igap = Inches(0.2)
    for i, (token, ok, hx, color) in enumerate(inks):
        x = MARGIN + (iw + igap) * i
        color_swatch(slide, x, y2, iw, h, token, ok, hx, color)
    # Urgency
    ux = MARGIN + (iw + igap) * 3 + Inches(0.4)
    add_text(slide, ux, Inches(4.85), Inches(4), Inches(0.3), "URGENCY ONLY",
             font=F_MONO, size=9, color=URGENCY, bold=True)
    color_swatch(slide, ux, y2, iw, h, "URGENCY",
                 "oklch(0.65 0.15 35)", "#d96847", URGENCY)


def make_typography(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "06   TYPOGRAPHY")
    slide_title(slide, "Three families. Nothing else.")
    # Three columns
    col_w = Inches(3.9)
    gap = Inches(0.3)
    y_top = Inches(2.8)
    families = [
        ("DISPLAY", "Bricolage", "Grotesque", "Headings, hero, tagline.\nWeight 600. Italic em is the punchline word."),
        ("UI", "Inter", "Tight", "Body copy, buttons, labels.\nWeight 400 body, 500 button, 600 emphasis."),
        ("MONO", "JetBrains", "Mono", "Eyebrows, labels, stats.\nUppercase, letter-spaced 0.18em, 10 to 11 px."),
    ]
    for i, (label, sample1, sample2, desc) in enumerate(families):
        x = MARGIN + (col_w + gap) * i
        # Label
        add_text(slide, x, y_top, col_w, Inches(0.3), label,
                 font=F_MONO, size=10, color=PURPLE_600, bold=True)
        # Big sample
        add_text(slide, x, y_top + Inches(0.45), col_w, Inches(0.9), sample1,
                 font=F_DISPLAY, size=38, color=INK, bold=True,
                 line_spacing=1.0)
        add_text(slide, x, y_top + Inches(1.2), col_w, Inches(0.9), sample2,
                 font=F_DISPLAY, size=38, color=INK, bold=True,
                 line_spacing=1.0)
        # Description
        add_text(slide, x, y_top + Inches(2.2), col_w, Inches(1.4), desc,
                 font=F_BODY, size=12, color=INK_3, line_spacing=1.4)


def make_type_scale(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "06   TYPE SCALE")
    slide_title(slide, "The scale, as it appears live.", size=36)
    # Fixed row layout: sample at top of row, label below, fixed gap between rows.
    samples = [
        ("Give Hope. Change a Life.",       "PAGE HERO H1   ·   600   ·   clamp(33, 9vw, 104)", 30, F_DISPLAY, True),
        ("Hopeful Survivors",                "SECTION H2   ·   600   ·   clamp(32, 4.4vw, 56)",  22, F_DISPLAY, True),
        ("A few of our survivors",           "SUB-SECTION H3   ·   600   ·   clamp(21, 2.5vw, 30)", 17, F_DISPLAY, True),
        ("Run 2 The Rescue pulls dogs from the meat trade.", "BODY LEAD   ·   400   ·   17 px", 13, F_BODY, False),
        ("DOGS RESCUED   ·   FORMS OF GIVING", "MONO LABEL   ·   600 uppercase   ·   10-11 px",     11, F_MONO,    True),
    ]
    y = Inches(3.1)
    for sample, label, size, font, bold in samples:
        sample_h = Inches(size / 60.0 + 0.15)  # space for the sample line
        add_text(slide, MARGIN, y, Inches(12), sample_h,
                 sample, font=font, size=size, color=INK, bold=bold,
                 line_spacing=1.0)
        add_text(slide, MARGIN, y + sample_h + Inches(0.02),
                 Inches(12), Inches(0.22),
                 label, font=F_MONO, size=8, color=INK_3)
        y += sample_h + Inches(0.46)


def make_photography(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "07   PHOTOGRAPHY")
    slide_title(slide, "Photography is the brand surface.")
    add_text(slide, MARGIN, Inches(3.0), Inches(11), Inches(0.5),
             "The right photo of a real survivor is worth a paragraph of copy.",
             font=F_BODY, size=16, color=INK_3, italic=True)
    # Use / Avoid columns
    col_w = Inches(5.8)
    use_x = MARGIN
    avoid_x = MARGIN + col_w + Inches(0.4)
    y_top = Inches(3.8)
    add_text(slide, use_x, y_top, col_w, Inches(0.3), "USE",
             font=F_MONO, size=10, color=PURPLE_700, bold=True)
    add_text(slide, avoid_x, y_top, col_w, Inches(0.3), "AVOID",
             font=F_MONO, size=10, color=URGENCY, bold=True)
    add_rect(slide, use_x, y_top + Inches(0.32), col_w, Inches(0.02), PURPLE_500)
    add_rect(slide, avoid_x, y_top + Inches(0.32), col_w, Inches(0.02), URGENCY)
    pairs = [
        ("Real R2TR dogs, named, with stories",  "Generic stock"),
        ("Eye level with the dog",                "Down-shots and clinical angles"),
        ("Soft natural light",                    "Harsh flash or studio glossy"),
        ("Quiet warmth",                          "Pity or exuberance"),
        ("Generous whitespace around portraits",  "Crowded composition"),
        ("Reality page (behind a notice)",        "Graphic cruelty anywhere else"),
    ]
    y = y_top + Inches(0.55)
    for use, avoid in pairs:
        add_text(slide, use_x, y, col_w, Inches(0.35), use,
                 font=F_BODY, size=14, color=INK, bold=True)
        add_text(slide, avoid_x, y, col_w, Inches(0.35), avoid,
                 font=F_BODY, size=14, color=INK_3, italic=True)
        y += Inches(0.45)


def make_iconography(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "08   ICONOGRAPHY")
    slide_title(slide, "Four marks. Used consistently.")
    items = [
        ("✦", "Sparkle",  "Leader before eyebrows. Purple 400 or 500. Never as body bullets."),
        ("✱", "Paw",      "Decorative. 0.08 to 0.20 opacity on dark, 0.15 to 0.30 on light. Mood, not message."),
        ("→", "Arrow",    "On hover, buttons nudge it 3 px right. Unicode →, not -> combo."),
        ("♡  ♥", "Hearts","Outline = favorite. Filled = confirmation. Always purple, never red."),
    ]
    y = Inches(2.7)
    for mark, name, desc in items:
        # Mark cell
        add_round_rect(slide, MARGIN, y, Inches(1.4), Inches(1.0),
                       PURPLE_SOFT, radius=0.1)
        add_text(slide, MARGIN, y, Inches(1.4), Inches(1.0),
                 mark, font=F_DISPLAY, size=40, color=PURPLE_700, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Name + desc
        add_text(slide, MARGIN + Inches(1.7), y + Inches(0.05),
                 Inches(10), Inches(0.4),
                 name, font=F_DISPLAY, size=20, color=INK, bold=True)
        add_text(slide, MARGIN + Inches(1.7), y + Inches(0.5),
                 Inches(10), Inches(0.5),
                 desc, font=F_BODY, size=13, color=INK_2)
        y += Inches(1.15)


def make_layout(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "09   LAYOUT AND WHITESPACE")
    slide_title(slide, "Whitespace is the brand.")
    add_text(slide, MARGIN, Inches(2.3), Inches(12), Inches(0.5),
             "More breathing room than the average nonprofit. On purpose.",
             font=F_BODY, size=16, color=INK_3, italic=True)
    rows = [
        ("SECTION PADDING",       "≥ 72 px desktop  ·  48 px mobile"),
        ("PAGE MARGINS",          "clamp(20px, 4vw, 48px)   ( --pad token )"),
        ("CARD PADDING",          "18 to 28 px"),
        ("HEADLINE → PARAGRAPH",  "12 to 22 px"),
        ("MAX CONTENT WIDTH",     "1200 px   ( --maxw token )"),
    ]
    y = Inches(3.4)
    for label, value in rows:
        add_text(slide, MARGIN, y, Inches(3.5), Inches(0.4), label,
                 font=F_MONO, size=10, color=PURPLE_600, bold=True)
        add_text(slide, MARGIN + Inches(3.7), y, Inches(8), Inches(0.4),
                 value, font=F_BODY, size=16, color=INK, bold=True)
        y += Inches(0.6)


def make_cta_stack(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PLUM_900)
    page_chrome(slide, page, total, dark=True)
    eyebrow(slide, MARGIN, Inches(0.8), "10   CTA PRIORITY STACK",
            color=PURPLE_400)
    slide_title(slide, "Four asks. One ladder.", color=WHITE)
    add_text(slide, MARGIN, Inches(2.3), Inches(12), Inches(0.5),
             "When CTAs compete for attention, this is the order.",
             font=F_BODY, size=16, color=NEAR_WHITE, italic=True)
    stack = [
        ("01", "ADOPT",   "Highest-impact outcome."),
        ("02", "SPONSOR", "Recurring, deepest commitment."),
        ("03", "DONATE",  "One-time gift."),
        ("04", "FOSTER",  "Supports the pipeline."),
    ]
    y = Inches(3.4)
    for num, label, desc in stack:
        add_text(slide, MARGIN, y, Inches(1.0), Inches(0.6), num,
                 font=F_DISPLAY, size=44, color=PURPLE_400, bold=True)
        add_text(slide, MARGIN + Inches(1.2), y + Inches(0.05),
                 Inches(3.5), Inches(0.6),
                 label, font=F_DISPLAY, size=32, color=WHITE, bold=True)
        add_text(slide, MARGIN + Inches(5.0), y + Inches(0.15),
                 Inches(7), Inches(0.6),
                 desc, font=F_BODY, size=16, color=NEAR_WHITE)
        y += Inches(0.85)


def make_cta_phrasing(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "10   CTA PHRASING")
    slide_title(slide, "Approved language.")
    col_w = Inches(5.8)
    use_x = MARGIN
    avoid_x = MARGIN + col_w + Inches(0.4)
    y_top = Inches(2.5)
    add_text(slide, use_x, y_top, col_w, Inches(0.3), "USE",
             font=F_MONO, size=10, color=PURPLE_700, bold=True)
    add_text(slide, avoid_x, y_top, col_w, Inches(0.3), "AVOID",
             font=F_MONO, size=10, color=URGENCY, bold=True)
    add_rect(slide, use_x, y_top + Inches(0.32), col_w, Inches(0.02), PURPLE_500)
    add_rect(slide, avoid_x, y_top + Inches(0.32), col_w, Inches(0.02), URGENCY)
    pairs = [
        ("Adopt today",                       "Adopt now!"),
        ("Meet the survivors",                "Save a life"),
        ("Bring one home",                    "Get a dog"),
        ("Start your monthly sponsorship",    "Sponsor us"),
        ("Be a Sponsor Angel",                "Give monthly"),
        ("Choose how to give",                "Please donate"),
        ("Apply to foster",                   "Sign up to foster"),
        ("Open your home",                    "Be a foster parent"),
    ]
    y = y_top + Inches(0.55)
    for use, avoid in pairs:
        add_text(slide, use_x, y, col_w, Inches(0.35), use,
                 font=F_BODY, size=14, color=INK, bold=True)
        add_text(slide, avoid_x, y, col_w, Inches(0.35), avoid,
                 font=F_BODY, size=14, color=INK_3, italic=True)
        y += Inches(0.45)


def make_applications(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "11   SAMPLE APPLICATIONS")
    slide_title(slide, "On every surface.")
    # Email signature card (top left)
    card_y = Inches(2.7)
    card_h = Inches(2.0)
    col1_w = Inches(5.8)
    col2_w = Inches(5.8)
    col2_x = MARGIN + col1_w + Inches(0.4)
    # Email card
    add_round_rect(slide, MARGIN, card_y, col1_w, card_h, WHITE)
    add_text(slide, MARGIN + Inches(0.25), card_y + Inches(0.2),
             col1_w, Inches(0.3), "EMAIL SIGNATURE",
             font=F_MONO, size=9, color=PURPLE_600, bold=True)
    add_paragraphs(slide,
        MARGIN + Inches(0.25), card_y + Inches(0.55), col1_w - Inches(0.5), Inches(1.4),
        [
            ("[Name]", {"size": 14, "bold": True, "color": INK}),
            ("Run 2 The Rescue   ·   501(c)(3) nonprofit", {"size": 12, "color": INK_2}),
            ("RUN. RESCUE. REPEAT.", {"size": 12, "color": PURPLE_700, "bold": True}),
            ("run2therescue.org", {"size": 12, "color": INK_3, "font": F_MONO}),
        ],
        font=F_BODY, line_spacing=1.4, space_after_pt=2)
    # Newsletter subject card
    add_round_rect(slide, col2_x, card_y, col2_w, card_h, WHITE)
    add_text(slide, col2_x + Inches(0.25), card_y + Inches(0.2),
             col2_w, Inches(0.3), "NEWSLETTER SUBJECT LINE",
             font=F_MONO, size=9, color=PURPLE_600, bold=True)
    add_text(slide, col2_x + Inches(0.25), card_y + Inches(0.7),
             col2_w - Inches(0.5), Inches(1.0),
             "Three new arrivals at LAX. Meet them.",
             font=F_DISPLAY, size=20, color=INK, bold=True, line_spacing=1.3)
    # Social caption (full width)
    sc_y = card_y + card_h + Inches(0.3)
    add_round_rect(slide, MARGIN, sc_y, col1_w + col2_w + Inches(0.4),
                   Inches(1.7), WHITE)
    add_text(slide, MARGIN + Inches(0.25), sc_y + Inches(0.2),
             Inches(8), Inches(0.3), "SOCIAL CAPTION",
             font=F_MONO, size=9, color=PURPLE_600, bold=True)
    add_paragraphs(slide,
        MARGIN + Inches(0.25), sc_y + Inches(0.55),
        col1_w + col2_w + Inches(0.0), Inches(1.2),
        [
            ('Meet [Dog Name]. Pulled from a holding pen in [Yulin / Seoul]. Cleared. Crated. And now waiting at JFK for a family.',
                {"size": 13, "color": INK, "space_after": 6}),
            ("Could it be yours? Adoption applications open. Link in bio.",
                {"size": 13, "color": INK_2, "space_after": 4}),
            ("#Run2TheRescue", {"size": 13, "color": PURPLE_600, "bold": True}),
        ],
        font=F_BODY, line_spacing=1.4)


def make_trust(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "12   TRUST SIGNALS")
    slide_title(slide, "Surface, don't bury.")
    rows = [
        ("EIN 99-4240461",                "Footer, donate page, press header."),
        ("501(c)(3) status",              "Phrasing on Donate and Footer."),
        ("Founders",                      "Cherven (CEO) and Klapper (COO) on Leadership."),
        ("Press",                         "Dodo, People, LI Press, NY Post. Linked."),
        ("IRS Determination Letter",      "Footer-linked PDF."),
        ("DAF legal name",                "Run to the Rescue (no 2). Use on DAF info card only."),
        ("14 years on the ground",        "Adopt hero stats."),
        ("Sources cited",                 "Mission stats credit HSI."),
    ]
    y = Inches(2.5)
    for label, note in rows:
        add_text(slide, MARGIN, y, Inches(4.5), Inches(0.4),
                 label, font=F_DISPLAY, size=16, color=INK, bold=True)
        add_text(slide, MARGIN + Inches(4.6), y, Inches(7.5), Inches(0.4),
                 note, font=F_BODY, size=13, color=INK_3, italic=True)
        y += Inches(0.5)


def make_retire(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LAV_50)
    page_chrome(slide, page, total)
    eyebrow(slide, MARGIN, Inches(0.8), "15   WHAT TO RETIRE")
    slide_title(slide, "Edit it the moment you see it.")
    col_w = Inches(5.8)
    use_x = MARGIN
    avoid_x = MARGIN + col_w + Inches(0.4)
    y_top = Inches(3.3)
    add_text(slide, use_x, y_top, col_w, Inches(0.3), "WRITE THIS",
             font=F_MONO, size=10, color=PURPLE_700, bold=True)
    add_text(slide, avoid_x, y_top, col_w, Inches(0.3), "NOT THIS",
             font=F_MONO, size=10, color=URGENCY, bold=True)
    add_rect(slide, use_x, y_top + Inches(0.32), col_w, Inches(0.02), PURPLE_500)
    add_rect(slide, avoid_x, y_top + Inches(0.32), col_w, Inches(0.02), URGENCY)
    pairs = [
        ("period or comma",                                  "em-dash (—) in body copy"),
        ("tax deductible",                                   "tax-deductible"),
        ("show what we believe by what we do",               "We believe …"),
        ("[cut it]",                                          "compassionate (as self-description)"),
        ("name the difference",                              "make a difference"),
        ("doors",                                            "asks"),
        ("real survivor photography",                        "stock dog photography"),
        ("hope",                                             "pity"),
        ("cage  ·  holding pen  ·  meat trade",              "slaughterhouse (in visible copy)"),
    ]
    y = y_top + Inches(0.55)
    for use, avoid in pairs:
        add_text(slide, use_x, y, col_w, Inches(0.35), use,
                 font=F_BODY, size=13, color=INK, bold=True)
        add_text(slide, avoid_x, y, col_w, Inches(0.35), avoid,
                 font=F_BODY, size=13, color=INK_3, italic=True)
        y += Inches(0.4)


def make_closing(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PLUM_900)
    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, PURPLE_500)
    page_chrome(slide, page, total, dark=True)
    eyebrow(slide, MARGIN, Inches(1.0), "ONE MORE THING", color=PURPLE_400)
    add_text(slide, MARGIN, Inches(2.0), Inches(12), Inches(2.6),
             "A playbook is a starting line.",
             font=F_DISPLAY, size=60, color=WHITE, bold=True,
             line_spacing=1.05)
    add_text(slide, MARGIN, Inches(4.8), Inches(12), Inches(0.5),
             "If a sentence is better because it breaks a rule here,",
             font=F_BODY, size=18, color=NEAR_WHITE)
    add_text(slide, MARGIN, Inches(5.2), Inches(12), Inches(0.5),
             "write the sentence. Ship the better thing.",
             font=F_BODY, size=18, color=NEAR_WHITE)
    # Divider
    add_rect(slide, MARGIN, Inches(6.15), Inches(3), Inches(0.03), PURPLE_500)
    add_text(slide, MARGIN, Inches(6.3), Inches(12), Inches(0.5),
             "RUN.   RESCUE.   REPEAT.",
             font=F_DISPLAY, size=20, color=PURPLE_400, bold=True)
    add_text(slide, MARGIN, Inches(6.75), Inches(12), Inches(0.4),
             '"Run to the rescue with love, and peace will follow."   — River Phoenix',
             font=F_BODY, size=11, color=NEAR_WHITE, italic=True)


# ---- Build ----

def build(out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Schedule slides
    builders = [
        make_cover,
        make_foreword,
        make_who_we_are,
        make_voice_attributes,
        make_reader_posture,
        make_transformation,
        make_protected_lexicon,
        make_color_plum,
        make_color_purple,
        make_color_lav_ink,
        make_typography,
        make_type_scale,
        make_photography,
        make_iconography,
        make_layout,
        make_cta_stack,
        make_cta_phrasing,
        make_applications,
        make_trust,
        make_retire,
        make_closing,
    ]
    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        if i == 1:
            fn(prs)
        else:
            fn(prs, i, total)

    prs.save(out_path)
    print(f"wrote {out_path} ({total} slides)")


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))
    out = os.path.join(here, "RUN_2_THE_RESCUE_BRAND_PLAYBOOK.pptx")
    build(out)
